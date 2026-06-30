"""
generate_presentation.py — Dynamic PPTX + slides JSON generator for OmniClient.

Public API
----------
create_presentation(topic, template, slide_count, slides_data) -> dict
    Returns: {"file_path": str, "slides_json": list[dict]}

build_slides_from_ai_response(ai_text, topic, slide_count) -> list[dict]
    Parse structured slide JSON from an AI chat response.

TEMPLATE_COLORS — dict of template name -> color constants
"""

from __future__ import annotations

import json
import os
import re
import uuid
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Directory where generated .pptx files are stored
PRESENTATIONS_DIR = Path(__file__).parent / "presentations"
PRESENTATIONS_DIR.mkdir(exist_ok=True)

# ---------------------------------------------------------------------------
# Template colour palettes
# ---------------------------------------------------------------------------

TEMPLATE_COLORS = {
    "Editorial": {
        "bg":       RGBColor(30,  41,  59),    # Slate 800
        "accent":   RGBColor(241, 245, 249),   # Slate 100
        "card":     RGBColor(51,  65,  85),    # Slate 700
        "text":     RGBColor(255, 255, 255),
        "muted":    RGBColor(203, 213, 225),
    },
    "Pixel": {
        "bg":       RGBColor(244, 63,  94),    # Rose 500
        "accent":   RGBColor(16,  185, 129),   # Emerald 500
        "card":     RGBColor(190, 24,  74),    # Rose 700
        "text":     RGBColor(255, 255, 255),
        "muted":    RGBColor(254, 205, 211),
    },
    "Vellum": {
        "bg":       RGBColor(217, 119, 6),     # Amber 600
        "accent":   RGBColor(254, 243, 199),   # Amber 100
        "card":     RGBColor(180, 83,  9),     # Amber 800
        "text":     RGBColor(255, 255, 255),
        "muted":    RGBColor(253, 230, 138),
    },
    "Sketch": {
        "bg":       RGBColor(2,   132, 199),   # Sky 600
        "accent":   RGBColor(224, 242, 254),   # Sky 100
        "card":     RGBColor(3,   105, 161),   # Sky 700
        "text":     RGBColor(255, 255, 255),
        "muted":    RGBColor(186, 230, 253),
    },
    "Whiteboard": {
        "bg":       RGBColor(17,  24,  39),    # Gray 900
        "accent":   RGBColor(249, 250, 251),   # Gray 50
        "card":     RGBColor(31,  41,  55),    # Gray 800
        "text":     RGBColor(255, 255, 255),
        "muted":    RGBColor(156, 163, 175),
    },
    "Minimal": {
        "bg":       RGBColor(243, 244, 246),   # Gray 100
        "accent":   RGBColor(75,  85,  99),    # Gray 600
        "card":     RGBColor(255, 255, 255),
        "text":     RGBColor(17,  24,  39),    # Gray 900
        "muted":    RGBColor(107, 114, 128),
    },
    "Corporate": {
        "bg":       RGBColor(30,  58,  138),   # Blue 900
        "accent":   RGBColor(59,  130, 246),   # Blue 500
        "card":     RGBColor(23,  37,  84),    # Blue 950
        "text":     RGBColor(255, 255, 255),
        "muted":    RGBColor(147, 197, 253),
    },
}
DEFAULT_TEMPLATE = "Corporate"


def _colors(template: str) -> dict:
    return TEMPLATE_COLORS.get(template, TEMPLATE_COLORS[DEFAULT_TEMPLATE])


# ---------------------------------------------------------------------------
# Slide data helpers
# ---------------------------------------------------------------------------

def _default_slides(topic: str, slide_count: int) -> list[dict]:
    """Generate a generic slide structure when no AI data is provided."""
    title = topic.strip().title()
    slides = [
        {
            "slide_number": 1,
            "type": "title",
            "title": title,
            "subtitle": f"A comprehensive overview of {title}",
            "notes": f"Welcome to this presentation on {title}.",
        },
        {
            "slide_number": 2,
            "type": "agenda",
            "title": "Agenda",
            "points": ["Introduction", "Key Concepts", "Analysis", "Recommendations", "Next Steps"],
            "notes": "Overview of what will be covered in this presentation.",
        },
    ]
    content_templates = [
        {"type": "content", "title": "Key Concepts", "points": [
            f"Core principle 1 of {title}",
            f"Core principle 2 of {title}",
            f"Core principle 3 of {title}",
        ]},
        {"type": "content", "title": "Analysis & Insights", "points": [
            "Data-driven insights",
            "Market trends and patterns",
            "Competitive landscape overview",
        ]},
        {"type": "content", "title": "Strategy", "points": [
            "Short-term priorities",
            "Medium-term goals",
            "Long-term vision",
        ]},
        {"type": "content", "title": "Implementation", "points": [
            "Phase 1: Foundation",
            "Phase 2: Execution",
            "Phase 3: Optimization",
        ]},
        {"type": "content", "title": "Results & KPIs", "points": [
            "Key performance indicators",
            "Success metrics",
            "Tracking and reporting",
        ]},
        {"type": "content", "title": "Case Studies", "points": [
            "Real-world example 1",
            "Real-world example 2",
            "Lessons learned",
        ]},
        {"type": "content", "title": "Challenges & Solutions", "points": [
            "Common pitfalls",
            "Mitigation strategies",
            "Best practices",
        ]},
        {"type": "content", "title": "Tools & Resources", "points": [
            "Essential tools",
            "Recommended resources",
            "Expert guidance",
        ]},
    ]
    for i in range(min(slide_count - 3, len(content_templates))):
        s = content_templates[i].copy()
        s["slide_number"] = i + 3
        s["notes"] = f"Detailed discussion of {s['title'].lower()}."
        slides.append(s)
    slides.append({
        "slide_number": len(slides) + 1,
        "type": "cta",
        "title": "Thank You",
        "subtitle": "Questions & Next Steps",
        "notes": "Open the floor for questions and outline immediate next steps.",
    })
    return slides[:slide_count]


def build_slides_from_ai_response(ai_text: str, topic: str, slide_count: int) -> list[dict]:
    """
    Extract structured slides JSON from an AI response text.
    Looks for a JSON code block first, then falls back to _default_slides.
    """
    # Try ```json ... ``` block
    m = re.search(r"```(?:json)?\s*(\[[\s\S]*?\])\s*```", ai_text, re.DOTALL)
    if m:
        try:
            data = json.loads(m.group(1))
            if isinstance(data, list) and data:
                return data[:slide_count]
        except json.JSONDecodeError:
            pass
    # Try raw JSON array
    m2 = re.search(r"(\[\s*\{[\s\S]*?\}\s*\])", ai_text, re.DOTALL)
    if m2:
        try:
            data = json.loads(m2.group(1))
            if isinstance(data, list) and data:
                return data[:slide_count]
        except json.JSONDecodeError:
            pass
    return _default_slides(topic, slide_count)


# ---------------------------------------------------------------------------
# PPTX rendering
# ---------------------------------------------------------------------------

def _apply_bg(slide, colors: dict, prs_width=13.333, prs_height=7.5):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(prs_width), Inches(prs_height),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors["bg"]
    bg.line.fill.background()
    # Accent bottom stripe
    acc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2),
        Inches(prs_width), Inches(0.3),
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = colors["accent"]
    acc.line.fill.background()


def _add_text(slide, left, top, width, height, text, font_size, bold=False, color=None,
              align=PP_ALIGN.LEFT, font_name="Inter", word_wrap=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return box


def _render_title_slide(slide, data: dict, colors: dict):
    _apply_bg(slide, colors)
    _add_text(slide, 1.0, 1.8, 11.333, 2.0,
              data.get("title", ""), 52, bold=True, color=colors["text"],
              align=PP_ALIGN.CENTER, font_name="Poppins")
    _add_text(slide, 1.5, 4.0, 10.333, 1.5,
              data.get("subtitle", ""), 22, bold=False, color=colors["muted"],
              align=PP_ALIGN.CENTER)


def _render_content_slide(slide, data: dict, colors: dict):
    _apply_bg(slide, colors)
    _add_text(slide, 1.0, 0.6, 11.333, 1.0,
              data.get("title", ""), 34, bold=True, color=colors["text"], font_name="Poppins")
    points = data.get("points") or data.get("bullets") or []
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.9),
        Inches(11.333), Inches(4.8),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = colors["card"]
    card.line.color.rgb = colors["accent"]
    card.line.width = Pt(1.5)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.35)
    first = True
    for pt in points:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "• " + str(pt)
        p.font.name = "Inter"
        p.font.size = Pt(16)
        p.font.color.rgb = colors["text"]
        p.space_before = Pt(10)


def _render_agenda_slide(slide, data: dict, colors: dict):
    _render_content_slide(slide, data, colors)


def _render_cta_slide(slide, data: dict, colors: dict):
    _apply_bg(slide, colors)
    _add_text(slide, 1.0, 2.0, 11.333, 1.8,
              data.get("title", "Thank You"), 52, bold=True, color=colors["text"],
              align=PP_ALIGN.CENTER, font_name="Poppins")
    if data.get("subtitle"):
        _add_text(slide, 1.5, 4.2, 10.333, 1.2,
                  data["subtitle"], 22, color=colors["muted"], align=PP_ALIGN.CENTER)


def _render_slide(prs, slide_data: dict, colors: dict):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide_type = slide_data.get("type", "content").lower()
    if slide_type == "title":
        _render_title_slide(slide, slide_data, colors)
    elif slide_type in ("agenda", "overview"):
        _render_agenda_slide(slide, slide_data, colors)
    elif slide_type in ("cta", "closing", "thank_you", "contact"):
        _render_cta_slide(slide, slide_data, colors)
    else:
        _render_content_slide(slide, slide_data, colors)
    # Speaker notes
    notes = slide_data.get("notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def create_presentation(
    topic: str = "Presentation",
    template: str = "Bold Blue",
    slide_count: int = 10,
    slides_data: Optional[list] = None,
) -> dict:
    """
    Generate a .pptx file and return metadata.

    Parameters
    ----------
    topic       : Human-readable topic/title
    template    : Template name from TEMPLATE_COLORS
    slide_count : Number of slides (clamped 8–20)
    slides_data : Optional list of slide dicts from AI; falls back to defaults

    Returns
    -------
    {
        "file_path":   str,       # Relative path under presentations/
        "slides_json": list[dict] # Slide data used
    }
    """
    slide_count = max(8, min(20, slide_count))
    colors = _colors(template)

    if not slides_data:
        slides_data = _default_slides(topic, slide_count)

    # Build PPTX
    prs_obj = Presentation()
    prs_obj.slide_width  = Inches(13.333)
    prs_obj.slide_height = Inches(7.5)

    for sd in slides_data:
        _render_slide(prs_obj, sd, colors)

    # Save to presentations dir
    safe_name = re.sub(r"[^\w\-]", "_", topic.strip().lower())[:40]
    filename   = f"{safe_name}_{uuid.uuid4().hex[:8]}.pptx"
    file_path  = PRESENTATIONS_DIR / filename
    prs_obj.save(str(file_path))

    return {
        "file_path":   str(file_path),
        "slides_json": slides_data,
    }


def build_slide_preview_html(slide: dict, template: str) -> str:
    """Generate a mini responsive HTML preview block for the active slide template."""
    bg_style = "background: linear-gradient(135deg, #1e293b, #0f172a);"
    text_color = "color: white;"
    card_bg = "background: rgba(255,255,255,0.08);"
    accent_bar = "background: #5f5af6;"
    
    if template == "Pixel":
        bg_style = "background: linear-gradient(135deg, #f43f5e, #be123c);"
        accent_bar = "background: #10b981;"
    elif template == "Vellum":
        bg_style = "background: linear-gradient(135deg, #d97706, #b45309);"
        accent_bar = "background: #fef3c7;"
    elif template == "Sketch":
        bg_style = "background: linear-gradient(135deg, #0284c7, #0369a1);"
        accent_bar = "background: #e0f2fe;"
    elif template == "Whiteboard":
        bg_style = "background: linear-gradient(135deg, #111827, #1f2937);"
        accent_bar = "background: #f9fafb;"
    elif template == "Minimal":
        bg_style = "background: linear-gradient(135deg, #f3f4f6, #e5e7eb);"
        text_color = "color: #111827;"
        card_bg = "background: white; border: 1px solid #d1d5db;"
        accent_bar = "background: #4b5563;"
    elif template == "Corporate":
        bg_style = "background: linear-gradient(135deg, #1e3a8a, #172554);"
        accent_bar = "background: #3b82f6;"

    title = slide.get("title", "")
    type_ = slide.get("type", "content").lower()
    
    html = f'<div class="slide-preview-box rounded-xl p-6 relative overflow-hidden h-48 flex flex-col justify-between" style="{bg_style} {text_color}">'
    html += f'<div class="absolute bottom-0 left-0 right-0 h-1" style="{accent_bar}"></div>'
    
    if type_ == "title":
        html += '<div class="text-center my-auto">'
        html += f'<h4 class="text-xl font-bold tracking-tight">{title}</h4>'
        html += f'<p class="text-xs opacity-80 mt-1">{slide.get("subtitle", "")}</p>'
        html += '</div>'
    elif type_ == "cta":
        html += '<div class="text-center my-auto">'
        html += f'<h4 class="text-xl font-bold tracking-tight">{title}</h4>'
        html += f'<p class="text-xs opacity-80 mt-1">{slide.get("subtitle", "")}</p>'
        html += '</div>'
    else:
        html += f'<h4 class="text-sm font-bold tracking-wide border-b border-white/10 pb-1 mb-2">{title}</h4>'
        html += f'<div class="p-3 rounded-lg flex-1 overflow-hidden" style="{card_bg}">'
        html += '<ul class="space-y-1">'
        points = slide.get("points") or slide.get("bullets") or []
        for pt in points[:3]:
            html += f'<li class="text-[10px] leading-tight list-disc list-inside truncate">{pt}</li>'
        html += '</ul>'
        html += '</div>'
        
    html += '</div>'
    return html


def generate_presentation_stream(
    topic: str,
    template: str,
    slide_count: int,
    conversation_id: Optional[int],
    db: Session,
):
    """
    Multi-phase presentation generator that yields SSE text events.
    1. researching: Web searches via search.py
    2. outlining: Outline generation via OpenAI
    3. generating: Detailed per-slide generation via OpenAI
    4. done: Save to DB & pptx build
    """
    from search import deep_search, format_search_for_context
    from config import get_settings
    from agent_engine import _make_client, _reasoning_kwargs
    from models import Presentation as PresentationModel
    
    settings = get_settings()
    client = _make_client()
    
    # ── PHASE 1: RESEARCHING ──
    q1 = f"{topic} latest trends key developments"
    q2 = f"{topic} industry analysis statistics"
    
    yield json.dumps({"type": "progress", "phase": "researching", "query": q1})
    try:
        results1 = deep_search(q1, db, max_results=3)
    except Exception:
        results1 = []
        
    yield json.dumps({"type": "progress", "phase": "researching", "query": q2})
    try:
        results2 = deep_search(q2, db, max_results=3)
    except Exception:
        results2 = []
        
    search_context = format_search_for_context(results1) + "\n\n" + format_search_for_context(results2)
    yield json.dumps({
        "type": "progress",
        "phase": "researching",
        "status": "completed",
        "summary": (search_context[:400] + "...") if search_context else "No results found."
    })
    
    # ── PHASE 2: OUTLINING ──
    outline_prompt = f"""You are a professional presentation architect. Based on this research context:
{search_context}

Generate an outline of exactly {slide_count} slides for the topic "{topic}".
Return a JSON array of slide objects. Each object must have keys:
- "slide_number" (int, 1-indexed)
- "title" (string, max 8 words)
- "description" (string, max 20 words summary of what this slide will cover)
- "type" (one of: "title", "agenda", "content", "cta")

Return ONLY the raw JSON array. No markdown code blocks, no other text."""

    try:
        resp = client.chat.completions.create(
            model=settings.default_model,
            messages=[{"role": "user", "content": outline_prompt}],
            temperature=0.3,
            **_reasoning_kwargs(),
        )
        outline_text = resp.choices[0].message.content.strip()
        if outline_text.startswith("```"):
            outline_text = re.sub(r"```(?:json)?\s*|```$", "", outline_text)
        outline = json.loads(outline_text)
    except Exception as e:
        # Fallback outline
        outline = [
            {"slide_number": 1, "title": topic.title(), "description": "Title slide", "type": "title"},
            {"slide_number": 2, "title": "Agenda", "description": "Presentation structure", "type": "agenda"},
        ]
        for i in range(3, slide_count):
            outline.append({"slide_number": i, "title": f"Key Topic {i-2}", "description": "Core information", "type": "content"})
        outline.append({"slide_number": slide_count, "title": "Conclusion", "description": "Closing remarks", "type": "cta"})

    yield json.dumps({"type": "progress", "phase": "outlining", "slides": outline})
    
    # ── PHASE 3: GENERATING ──
    slides_data = []
    for idx, item in enumerate(outline):
        slide_num = item["slide_number"]
        generating_prompt = f"""You are a professional presentation writer. Based on the research:
{search_context}

We are writing slide {slide_num} of {slide_count}.
Slide outline title: "{item['title']}"
Slide description: "{item['description']}"
Slide type: "{item['type']}"

Generate the full slide content as a JSON object with:
- "slide_number": {slide_num}
- "type": "{item['type']}"
- "title": "{item['title']}"
- "subtitle": string (ONLY if type is "title" or "cta", otherwise empty)
- "points": list of 3-4 concise bullet points (ONLY if type is "content" or "agenda", otherwise empty)
- "notes": detailed speaker notes (2-3 sentences)

CRITICAL RULES:
1. Use only data gathered from the research phase. If exact figures aren't available, clearly label estimates as estimates. Never invent precise-looking fake statistics.
2. Return ONLY the raw JSON object, no explanation, no markdown blocks."""

        try:
            resp = client.chat.completions.create(
                model=settings.default_model,
                messages=[{"role": "user", "content": generating_prompt}],
                temperature=0.4,
                **_reasoning_kwargs(),
            )
            slide_text = resp.choices[0].message.content.strip()
            if slide_text.startswith("```"):
                slide_text = re.sub(r"```(?:json)?\s*|```$", "", slide_text)
            slide_obj = json.loads(slide_text)
        except Exception:
            # Fallback slide
            slide_obj = {
                "slide_number": slide_num,
                "type": item["type"],
                "title": item["title"],
                "subtitle": f"Discussion about {topic}" if item["type"] in ("title", "cta") else "",
                "points": [f"Key point 1 regarding {item['title']}", f"Supporting details for {topic}"],
                "notes": f"This slide presents insights on {item['title']}."
            }
            
        slides_data.append(slide_obj)
        slide_html = build_slide_preview_html(slide_obj, template)
        
        yield json.dumps({
            "type": "progress",
            "phase": "generating",
            "slide_number": slide_num,
            "total": slide_count,
            "slide_html": slide_html,
            "slide_data": slide_obj
        })

    # ── PHASE 4: DELIVER (DONE) ──
    try:
        result = create_presentation(
            topic=topic,
            template=template,
            slide_count=slide_count,
            slides_data=slides_data,
        )
        file_path = result["file_path"]
    except Exception:
        file_path = ""
        
    pres = PresentationModel(
        title=topic.strip().title(),
        topic=topic,
        template=template,
        slide_count=slide_count,
        file_path=file_path,
        slides_json=json.dumps(slides_data),
        status="completed",
        conversation_id=conversation_id,
    )
    db.add(pres)
    db.commit()
    db.refresh(pres)
    
    yield json.dumps({
        "type": "progress",
        "phase": "done",
        "id": pres.id,
        "project_id": pres.id,
        "file_path": file_path,
        "title": pres.title,
        "topic": pres.topic,
        "template": pres.template,
        "slide_count": pres.slide_count,
        "created_at": pres.created_at.isoformat() if pres.created_at else None,
        "slides": slides_data,
    })


if __name__ == "__main__":
    result = create_presentation(
        topic="Digital Marketing Strategy for SaaS",
        template="Bold Blue",
        slide_count=10,
    )
    print("Generated:", result["file_path"])
    print("Slides:", len(result["slides_json"]))



if __name__ == "__main__":
    result = create_presentation(
        topic="Digital Marketing Strategy for SaaS",
        template="Bold Blue",
        slide_count=10,
    )
    print("Generated:", result["file_path"])
    print("Slides:", len(result["slides_json"]))



